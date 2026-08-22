"""임상수기센터 협업 제안 덱(9장) 생성. 레퍼런스 덱과 같은 16:9·배색·구조, 글자 크기 +10%."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
import os
import glob

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(ROOT, '..', '..', 'docs', 'keimyung-proposal', 'assets')
OUT = os.path.join(ROOT, '..', '..', 'docs', 'keimyung-proposal', '렉처링크_CPX_임상수기센터_협업제안.pptx')

W, H = 24384000, 13716000          # 레퍼런스 덱과 동일
FONT = 'Apple SD Gothic Neo'       # Windows 에서는 맑은 고딕으로 대체됨
INK, INK2, INK3 = RGBColor(0x2A,0x2A,0x2A), RGBColor(0x51,0x41,0x49), RGBColor(0x7A,0x71,0x76)
GREEN, RED = RGBColor(0x1B,0x9E,0x3A), RGBColor(0xE0,0x31,0x2D)
MINT, MINT2, WHITE, CREAM = RGBColor(0xB5,0xEF,0xC6), RGBColor(0xD6,0xF5,0xDF), RGBColor(255,255,255), RGBColor(0xFF,0xFA,0xF7)
PH_FILL, PH_LINE, SOFT = RGBColor(0xF1,0xEC,0xE8), RGBColor(0xB9,0xB0,0xAC), RGBColor(0xE4,0xDC,0xD8)

X = lambda p: Emu(int(W * p / 100))
Y = lambda p: Emu(int(H * p / 100))

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(W), Emu(H)
BLANK = prs.slide_layouts[6]


def text(slide, x, y, w, h, runs, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line=1.25):
    """runs: str 또는 [(text, {color,bold,size}), ...]. '\n' 으로 줄바꿈."""
    tb = slide.shapes.add_textbox(X(x), Y(y), X(w), Y(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str): runs = [(runs, {})]
    paras = [[]]
    for t, st in runs:
        for i, part in enumerate(t.split('\n')):
            if i: paras.append([])
            if part: paras[-1].append((part, st))
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line
        for t, st in para:
            r = p.add_run(); r.text = t
            f = r.font; f.name = FONT; f.size = Pt(st.get('size', size)); f.bold = st.get('bold', bold)
            f.color.rgb = st.get('color', color)
    return tb


def rect(slide, x, y, w, h, fill, line=None, dash=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(x), Y(y), X(w), Y(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1.2)
        if dash: s.line.dash_style = MSO_LINE.DASH
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def hline(slide, x, y, w, color=INK2, weight=1.0):
    c = slide.shapes.add_connector(1, X(x), Y(y), X(x + w), Y(y))
    c.line.color.rgb = color; c.line.width = Pt(weight)


def placeholder(slide, x, y, w, h, label):
    rect(slide, x, y, w, h, PH_FILL, PH_LINE, dash=True)
    text(slide, x, y, w, h, label, 22, color=INK3, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.4)


def cover_bg(slide):
    # 레퍼런스 표지: 그라데이션 이미지를 좌측으로 넘치게 배치
    slide.shapes.add_picture(os.path.join(ASSETS, 'bg-cover.png'), X(-52), Y(0), X(171), Y(120))


def content_bg(slide):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = CREAM
    slide.shapes.add_picture(os.path.join(ASSETS, 'bg-band.png'), X(0), Y(0), X(100), Y(72))
    text(slide, 4, 4, 40, 4, 'LectureLink', 20, color=INK2)
    text(slide, 56, 4, 40, 4, '전재현, 장유림', 20, color=INK2, align=PP_ALIGN.RIGHT)
    hline(slide, 4, 8.6, 92)


def title_block(slide, title, sub_runs):
    # 레퍼런스 본문 제목 72pt → 80pt, 부제 50 → 55 (+10%)
    text(slide, 4, 11, 92, 11, title, 80, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 4, 21.5, 92, 7, sub_runs, 55, bold=True, color=INK2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def panel(slide, x=3.5, y=31, w=93, h=60):
    return rect(slide, x, y, w, h, WHITE)


BODY, FINE = 36, 28   # 32→36, 25→28 (+10%)
DATA = os.path.join(ROOT, '..', '..', 'services', 'cpx', 'data', 'cpx')
N_CASES = len(glob.glob(os.path.join(DATA, 'cases', '*', '*.json')))
N_CC = len(glob.glob(os.path.join(DATA, 'cases', '*')))

# ───────────────────────── S01 표지
s = prs.slides.add_slide(BLANK); cover_bg(s)
text(s, 7, 18, 70, 5, '계명대학교 의과대학 임상수기센터 협업 제안', 33, color=INK2)
text(s, 7, 24, 75, 14, 'LectureLink CPX', 120, bold=True)
text(s, 7, 38, 80, 16, 'AI 표준화 환자와 혼자 하는\n진료수행 연습', 62, bold=True, line=1.2)
hline(s, 7, 57, 51)
for i, (k, v) in enumerate([('제안', '렉처링크 개발팀 · 전재현, 장유림 (계명의대 의학과)'),
                            ('요청', 'CPX 콘텐츠 의학 검수 · 공동연구 · 정식 협약 논의'),
                            ('자료', '상세 문서 별첨 (인쇄본) · lecturelink.kro.kr')]):
    text(s, 7, 60 + i * 6.2, 9, 5, k, 30, bold=True)
    text(s, 16, 60 + i * 6.2, 70, 5, v, 30)
s.shapes.add_picture(os.path.join(ASSETS, 'logo.png'), X(75), Y(70), X(25), Y(30))

# ───────────────────────── S02 한 번의 연습 (노트북 시연의 대체 화면)
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '학생 한 명이 혼자, 12분 안에',
            [('실제 CPX와 같은 흐름 — ', {}), ('문진 · 신체진찰 · 환자교육', {'color': GREEN}), (' 후 즉시 채점', {})])
panel(s)
for i, lab in enumerate(['화면 캡처 ①\n음성 문진', '화면 캡처 ②\n부위별 신체진찰', '화면 캡처 ③\n채점 결과']):
    placeholder(s, 6 + i * 30, 35, 27.5, 52, lab)

# ───────────────────────── S03 시나리오 구성
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, f'국시원 공개 항목표의 주호소 {N_CC}개 전부',
            [('주호소마다 감별진단별 증례를 두어 ', {}), (f'{N_CASES}개 시나리오', {'color': GREEN})])
panel(s)
for i, (n, l) in enumerate([(str(N_CC), '주호소'), (str(N_CASES), '증례')]):
    text(s, 6 + i * 22, 42, 20, 16, n, 120, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6 + i * 22, 60, 20, 6, l, 34, bold=True, color=INK2, align=PP_ALIGN.CENTER)
text(s, 52, 38, 43, 50, [
    ('수면장애 16 · 두통 8 · 가슴 통증 7 · 급성 복통 6 …\n\n', {'bold': True, 'size': BODY}),
    ('예: 가슴 통증 →\n', {'size': BODY}),
    ('급성심근경색 · 대동맥박리 · 기흉 · 폐색전 ·\n안정협심증 · 역류성식도염 · 늑연골염', {'size': FINE, 'color': INK3}),
], BODY, color=INK2, anchor=MSO_ANCHOR.MIDDLE, line=1.35)

# ───────────────────────── S04 채점표
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '채점표는 실제 CPX 배점 구조를 따릅니다',
            [('체크리스트 70 + 환자의사관계 30 · 항목마다 ', {}), ('대화 근거', {'color': GREEN}), ('를 붙여 판정', {})])
panel(s)
for i, (name, v) in enumerate([('병력청취', 38), ('신체진찰', 16), ('환자교육', 16), ('환자의사관계', 30)]):
    y = 37 + i * 12.5
    text(s, 6, y, 20, 5, name, 30, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 6, y + 5.5, 34, 3.2, MINT2)
    rect(s, 6, y + 5.5, 34 * v / 40, 3.2, GREEN)
    text(s, 41, y, 6, 5, str(v), 34, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, 53, 36, 42, 52, [
    ('판정은 ', {'size': BODY}), ('충족 · 부분 · 미충족', {'size': BODY, 'bold': True}), (' 세 단계\n\n', {'size': BODY}),
    ('“청진을 해보겠습니다” → 시행으로 인정\n', {'size': FINE, 'color': INK3}),
    ('단어가 달라도 문맥이 같으면 인정\n', {'size': FINE, 'color': INK3}),
    ('점수 계산은 AI가 아니라 정해진 규칙이 합니다', {'size': FINE, 'color': INK3}),
], BODY, color=INK2, anchor=MSO_ANCHOR.MIDDLE, line=1.4)

# ───────────────────────── S05 표준화 환자
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '환자는 묻는 것에만 답합니다',
            [('표준화 환자 교육 원칙을 그대로 ', {}), ('행동 규칙', {'color': GREEN}), ('으로', {})])
panel(s)
items = ['▶ 진단명·핵심 단서를 먼저 말하지 않음', '▶ 민감한 질문엔 이유를 들어야 협조',
         '▶ 진찰 요청 시 행동 반응 → 소견은 객관 자료로', '▶ 같은 증례라도 이름·성별·직업은 매번 달라짐']
for i, t in enumerate(items):
    text(s, 7 + (i % 2) * 46, 42 + (i // 2) * 20, 44, 12, t, BODY, color=INK2, anchor=MSO_ANCHOR.MIDDLE, line=1.3)

# ───────────────────────── S06 일관성
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '같은 수행에는 같은 점수가 나옵니다',
            [('동일 대화 기록 10회 재채점 — ', {}), ('일관성', {'color': GREEN}), ('은 확인했습니다', {})])
panel(s)
for i, (k, v) in enumerate([('항목 판정 일치', '32항목 중 31개 동일'), ('총점 편차', '±0.75점'), ('채점표 형식 오류', '0건')]):
    y = 40 + i * 13
    text(s, 7, y, 22, 8, k, 32, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 29, y, 24, 8, v, 32, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    hline(s, 7, y + 9.5, 46, SOFT, 0.75)
text(s, 58, 36, 37, 52, [('일관성', {'color': RED, 'bold': True}), ('입니다.\n', {}), ('정확성', {'bold': True}), ('은 아직 아닙니다.', {})],
     44, color=INK2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.4)

# ───────────────────────── S07 한계선 (장 구분 유형)
s = prs.slides.add_slide(BLANK); cover_bg(s)
text(s, 4, 5, 30, 4, 'LectureLink', 20, color=INK2)
text(s, 66, 5, 30, 4, '전재현, 장유림', 20, color=INK2, align=PP_ALIGN.RIGHT)
hline(s, 4, 9.2, 92)
text(s, 5, 30, 14, 40, '?', 220, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
text(s, 20, 33, 75, 14, '임상적으로 맞는가', 96, bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, 20, 48, 75, 20, [('▷ 검수를 마친 증례 ', {}), ('0', {'color': RED}), (f' / {N_CASES}\n', {}),
                         ('▷ 채점표는 공개 항목표 + 교재 참조 자체 제작', {})], 50, bold=True, color=INK2, line=1.3)

# ───────────────────────── S08 검수 요청
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '의학 검수를 부탁드립니다',
            [('도구와 표본은 저희가 준비하고, ', {}), ('판단', {'color': GREEN}), ('만 맡겨 주십시오', {})])
steps = [('1단계 · 2주', '채점표 정본 대조', '센터 채점표 1~2종과\n항목 단위 비교', True),
         ('2단계 · 4주', '증례 블라인드 검수', '층화 추출 표본,\n검수 화면 제공', False),
         ('3단계 · 6주', '학생 파일럿', '센터 실습과 병행,\n결과 공동 분석', False)]
for i, (k, t, d, on) in enumerate(steps):
    x = 4 + i * 31
    rect(s, x, 36, 30, 44, WHITE); rect(s, x, 36, 30, 1.6, GREEN if on else MINT)
    text(s, x + 2.5, 41, 25, 5, k, 24, color=INK3)
    text(s, x + 2.5, 47, 25, 10, t, 40, bold=True, line=1.2)
    text(s, x + 2.5, 60, 25, 18, d, FINE, color=INK2, line=1.35)

# ───────────────────────── S09 공동연구 제안
s = prs.slides.add_slide(BLANK); content_bg(s)
title_block(s, '함께 하고 싶은 연구',
            [('채점 타당도', {'color': GREEN}), ('를 먼저, 학습 효과는 그다음 — 설계와 일정은 함께 정합니다', {})])
cols = [('연구 ①  채점 타당도', 'AI 판정과 전문가 판정은 얼마나, 어디서 갈리는가',
         ['일치도는 어느 수준인가', '불일치는 어떤 항목에 몰리는가', '인용한 대화 근거가 실제와 맞는가', '재채점해도 같은 결과가 나오는가']),
        ('연구 ②  학습 효과', '플랫폼 밖의 준거로 확인하는 학습 효과',
         ['실습 평가 등 외부 준거로 향상되는가', '피드백 방식에 따라 달라지는가', '같은 증례 반복과 변형 중 무엇이 낫는가', '실습 과정에 통합할 수 있는가'])]
for i, (head, sub, qs) in enumerate(cols):
    x = 4 + i * 47
    rect(s, x, 32, 45, 46, WHITE); rect(s, x, 32, 45, 1.6, GREEN if i == 0 else MINT)
    text(s, x + 3, 36, 40, 6, head, 36, bold=True)
    text(s, x + 3, 42.5, 40, 6, sub, 26, color=INK3)
    text(s, x + 3, 50, 40, 32, '\n'.join('▷ ' + q for q in qs), FINE, color=INK2, line=1.45)
text(s, 4, 83, 92, 8, [('저희가 먼저 약속하는 것  ', {'bold': True, 'color': INK}),
                        ('부정적 결과도 발표에 동의 · 데이터와 분석 코드 공개 · 분석 계획 사전 등록', {})],
     26, color=INK2, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print('saved', OUT, len(prs.slides), 'slides')
